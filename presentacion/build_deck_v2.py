"""Construye la presentacion ejecutiva v2 sin modificar el PPTX original."""

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


BASE = Path(__file__).resolve().parent
ASSETS = BASE / "assets"
OUTPUT = BASE / "Presentacion_Ejecutiva_Flip_CABA_v4.pptx"

with (ASSETS / "datos_presentacion.json").open(encoding="utf-8") as fh:
    DATA = json.load(fh)
with (ASSETS / "datos_ejemplo.json").open(encoding="utf-8") as fh:
    EXAMPLE = json.load(fh)

NAVY = "0F1B2D"
TERRA = "C2683C"
PAPER = "F4F1EC"
WHITE = "FFFFFF"
INK = "3A3A3A"
GRAY = "7A7A7A"
LIGHT = "D8D0C2"
ICE = "C7D2DD"
BLUE = "4C78A8"
SERIF = "Georgia"
SANS = "Calibri"

W = 13.333
H = 7.5
MX = 0.7


def rgb(value):
    return RGBColor.from_string(value)


def add_rect(slide, x, y, w, h, color, line=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.color.rgb = rgb(line or color)
    return shape


def add_line(slide, x, y, w, color=LIGHT, height=0.014):
    return add_rect(slide, x, y, w, height, color)


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=14,
    color=INK,
    font=SANS,
    bold=False,
    italic=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)
    return box


def add_rich_text(slide, runs, x, y, w, h, size=14, color=INK, font=SANS):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    paragraph = frame.paragraphs[0]
    paragraph.space_after = Pt(0)
    for item in runs:
        run = paragraph.add_run()
        run.text = item["text"]
        run.font.name = item.get("font", font)
        run.font.size = Pt(item.get("size", size))
        run.font.bold = item.get("bold", False)
        run.font.italic = item.get("italic", False)
        run.font.color.rgb = rgb(item.get("color", color))
    return box


def add_bullets(slide, items, x, y, w, h, size=14, color=INK, spacing=8):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, item in enumerate(items):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = SANS
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(color)
        p.space_after = Pt(spacing)
        p.text = f"•  {item}"
    return box


def add_picture(slide, filename, x, y, w=None, h=None):
    kwargs = {}
    if w is not None:
        kwargs["width"] = Inches(w)
    if h is not None:
        kwargs["height"] = Inches(h)
    return slide.shapes.add_picture(
        str(ASSETS / filename), Inches(x), Inches(y), **kwargs
    )


def add_header(slide, section, page, title, dark=False, title_size=30):
    background = NAVY if dark else PAPER
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(background)
    add_text(slide, section.upper(), MX, 0.32, 9.2, 0.3, 10, TERRA, bold=True)
    add_text(
        slide,
        f"{page:02d} / 12",
        11.0,
        0.32,
        1.55,
        0.3,
        10,
        ICE if dark else GRAY,
        align=PP_ALIGN.RIGHT,
    )
    add_line(slide, MX, 0.72, W - 2 * MX, "2A3C52" if dark else LIGHT)
    add_text(
        slide,
        title,
        MX,
        0.87,
        W - 2 * MX,
        0.85,
        title_size,
        WHITE if dark else NAVY,
        SERIF,
        bold=True,
    )


def add_step(slide, number, title, body, x, y, w):
    add_text(slide, number, x, y, 0.7, 0.45, 22, TERRA, SERIF, bold=True)
    add_text(slide, title, x, y + 0.52, w, 0.45, 20, NAVY, SERIF, bold=True)
    add_line(slide, x, y + 1.05, w - 0.15)
    add_text(slide, body, x, y + 1.18, w - 0.1, 1.4, 12.5, INK)


prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
blank = prs.slide_layouts[6]
prs.core_properties.author = "Felipe Tamaki, Matias Goldschmidt, Lucas Binelo"
prs.core_properties.title = "Priorizacion de oportunidades de flip inmobiliario en CABA"
prs.core_properties.subject = "Presentacion ejecutiva"

# 1. Portada
slide = prs.slides.add_slide(blank)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = rgb(NAVY)
add_picture(slide, DATA["assets"]["portada"], 7.15, 0.55, h=6.35)
add_text(slide, "ANÁLISIS DE OPORTUNIDADES INMOBILIARIAS · CABA", 0.85, 0.58, 7, 0.3, 10.5, TERRA, bold=True)
add_text(slide, "ITBA · 2026", 9.3, 0.58, 3.1, 0.3, 10.5, ICE, align=PP_ALIGN.RIGHT)
add_line(slide, 0.85, 0.97, 11.65, "2A3C52")
add_text(slide, "Priorización de\noportunidades de\nflip inmobiliario", 0.82, 1.85, 6.4, 2.55, 40, WHITE, SERIF, bold=True)
add_text(slide, "Una herramienta para decidir dónde buscar, qué propiedades investigar y qué validar antes de invertir.", 0.85, 4.65, 5.8, 0.95, 15, ICE)
add_line(slide, 0.85, 5.95, 3.0, TERRA)
add_text(slide, "Felipe Tamaki · 66477   |   Matías Goldschmidt · 66061   |   Lucas Binelo · 66011", 0.85, 6.15, 7.3, 0.35, 11.5, WHITE, bold=True)
add_text(slide, "Analítica Descriptiva · Presentación ejecutiva", 0.85, 6.55, 5.5, 0.3, 10.5, ICE)

# 2. Problema
slide = prs.slides.add_slide(blank)
add_header(slide, "Contexto y oportunidad", 2, "Miles de avisos, pero poca evidencia para decidir")
add_text(slide, "El desafío no es encontrar avisos, sino compararlos antes de comprometer capital.", MX, 1.95, 6.2, 1.0, 21, NAVY, SERIF, bold=True)
add_text(slide, "La información es heterogénea y las señales sobre precio, estado y entorno están dispersas. Sin un proceso analítico, la selección depende demasiado de intuición y tiempo de búsqueda.", MX, 3.22, 6.15, 1.75, 15, INK)
stats = [
    ("12.518", "avisos relevados"),
    ("48", "barrios con dinámicas distintas"),
    ("USD 100k+", "capital típico comprometido"),
]
for i, (value, label) in enumerate(stats):
    y = 1.9 + i * 1.48
    if i:
        add_line(slide, 8.15, y - 0.16, 4.35)
    add_text(slide, value, 8.3, y, 4.0, 0.68, 37, NAVY, SERIF, bold=True)
    add_text(slide, label, 8.34, y + 0.72, 4.0, 0.35, 12, GRAY)
add_rect(slide, 8.03, 1.95, 0.02, 4.25, TERRA)

# 3. Modelo de negocio flipper
slide = prs.slides.add_slide(blank)
add_header(slide, "Modelo de negocio", 3, "El resultado de un flip se define antes de empezar la obra")
flipper_steps = [
    ("01", "Comprar bien", "Entrar por debajo de propiedades comparables."),
    ("02", "Mejorar con criterio", "Invertir solo en cambios que el mercado del barrio valore."),
    ("03", "Revender con margen", "Salir en un mercado con demanda y referencias suficientes."),
]
for i, (num, title, body) in enumerate(flipper_steps):
    x = MX + i * 4.13
    add_text(slide, num, x, 2.0, 0.65, 0.42, 21, TERRA, SERIF, bold=True)
    add_text(slide, title, x, 2.58, 3.55, 0.5, 22, NAVY, SERIF, bold=True)
    add_line(slide, x, 3.2, 3.45)
    add_text(slide, body, x, 3.43, 3.45, 1.0, 14, INK)
    if i < 2:
        add_text(slide, "→", x + 3.55, 3.45, 0.45, 0.45, 20, TERRA, align=PP_ALIGN.CENTER)
add_rect(slide, MX, 5.02, W - 2 * MX, 1.08, NAVY)
add_text(slide, "Margen potencial  =  venta esperada  −  compra negociada  −  obra  −  gastos  −  tiempo", 1.0, 5.3, 11.3, 0.45, 19, WHITE, SERIF, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, "Nuestro análisis mejora la selección de compra. Todavía no calcula rentabilidad porque faltan precios de cierre y costos reales.", MX, 6.42, 11.8, 0.52, 12.5, GRAY, italic=True, align=PP_ALIGN.CENTER)

# 4. Datos
slide = prs.slides.add_slide(blank)
add_header(slide, "Datos utilizados", 4, "Integramos publicaciones con información territorial")
sources = [
    ("ARGENPROP", "Publicaciones de departamentos en venta", "Precio, superficie, ambientes, estado y amenities"),
    ("GCBA", "Datos abiertos de la Ciudad", "Barrios, comunas, subte, hospitales y contexto urbano"),
    ("OPENSTREETMAP", "Servicios y geocodificación", "Colegios, supermercados, avenidas y transporte"),
]
for i, (name, role, detail) in enumerate(sources):
    x = MX + i * 4.12
    add_text(slide, f"0{i + 1}", x, 2.0, 0.65, 0.45, 22, TERRA, SERIF, bold=True)
    add_text(slide, name, x, 2.58, 3.55, 0.38, 15, NAVY, bold=True)
    add_line(slide, x, 3.08, 3.45)
    add_text(slide, role, x, 3.28, 3.45, 0.8, 15, INK, SERIF, bold=True)
    add_text(slide, detail, x, 4.25, 3.45, 1.0, 12.5, GRAY)
    if i < 2:
        add_text(slide, "+", x + 3.62, 3.35, 0.4, 0.5, 22, TERRA, SERIF, bold=True, align=PP_ALIGN.CENTER)
add_line(slide, MX, 5.75, W - 2 * MX, LIGHT)
add_text(slide, "Cada fila representa una publicación. El análisis trabaja con precios publicados, no con precios efectivos de cierre.", MX, 5.95, 11.8, 0.6, 14, NAVY, SERIF, italic=True)

# 5. Señales buscadas
slide = prs.slides.add_slide(blank)
add_header(slide, "Qué buscamos en los datos", 5, "Cuatro señales para decidir qué avisos merecen una revisión")
signals = [
    ("01", "Precio de entrada", "¿Publica por debajo de inmuebles del mismo barrio, tipo y ambientes?"),
    ("02", "Margen de mejora", "¿El estado permite una puesta en valor sin ocultar un problema estructural?"),
    ("03", "Mercado de salida", "¿Hay oferta suficiente para construir comparables y entender la zona?"),
    ("04", "Segunda opinión", "¿El modelo también detecta que el precio parece desalineado?"),
]
for i, (num, title, body) in enumerate(signals):
    x = MX + i * 3.08
    add_text(slide, num, x, 2.05, 0.6, 0.4, 20, TERRA, SERIF, bold=True)
    add_text(slide, title, x, 2.6, 2.7, 0.5, 19, NAVY, SERIF, bold=True)
    add_line(slide, x, 3.22, 2.72)
    add_text(slide, body, x, 3.48, 2.72, 1.55, 13, INK)
add_line(slide, MX, 5.55, W - 2 * MX)
add_text(slide, "12.518 avisos", MX, 5.84, 2.2, 0.5, 22, NAVY, SERIF, bold=True)
add_text(slide, "→", 2.72, 5.88, 0.55, 0.4, 18, TERRA, align=PP_ALIGN.CENTER)
add_text(slide, "7.245 propiedades comparables", 3.28, 5.84, 3.65, 0.5, 22, NAVY, SERIF, bold=True)
add_text(slide, "→", 7.05, 5.88, 0.55, 0.4, 18, TERRA, align=PP_ALIGN.CENTER)
add_text(slide, "3.129 casos para investigar", 7.62, 5.84, 3.9, 0.5, 22, TERRA, SERIF, bold=True)
add_text(slide, "Geocodificación, enriquecimiento y limpieza permiten comparar; el ranking reduce el universo.", MX, 6.42, 11.8, 0.4, 11.5, GRAY, italic=True)

# 6. Solución
slide = prs.slides.add_slide(blank)
add_header(slide, "La solución", 6, "Transformamos datos en un proceso de decisión")
flow = [
    ("01", "Explorar", "Entender precios, oferta y diferencias territoriales."),
    ("02", "Comparar", "Construir referencias entre inmuebles equivalentes."),
    ("03", "Priorizar", "Ordenar barrios y propiedades con señales explicables."),
    ("04", "Contrastar", "Usar el modelo como segunda opinión de precio."),
    ("05", "Validar", "Completar la due diligence antes de invertir."),
]
for i, (num, title, body) in enumerate(flow):
    x = MX + i * 2.45
    add_text(slide, num, x, 2.12, 0.6, 0.4, 20, TERRA, SERIF, bold=True)
    add_text(slide, title, x, 2.7, 2.15, 0.42, 18, NAVY, SERIF, bold=True)
    add_line(slide, x, 3.27, 2.05)
    add_text(slide, body, x, 3.52, 2.05, 1.55, 12, INK)
    if i < 4:
        add_text(slide, "→", x + 2.05, 4.0, 0.4, 0.4, 18, TERRA, align=PP_ALIGN.CENTER)
add_line(slide, MX, 5.65, W - 2 * MX)
add_text(slide, "Resultado: un embudo barrio → propiedad → validación. Reduce el universo; no automatiza la compra.", MX, 5.88, 11.8, 0.6, 15, NAVY, SERIF, italic=True)

# 7. Insight: ubicacion y producto
slide = prs.slides.add_slide(blank)
add_header(slide, "Insight 1 · Ubicación y producto", 7, "El barrio fija el punto de entrada; los amenities agregan valor, pero no alcanzan solos", title_size=25)
add_text(slide, "DÓNDE ESTÁ", MX, 1.66, 5.7, 0.3, 10, TERRA, bold=True)
add_text(slide, "Precio publicado por barrio", MX, 1.96, 5.7, 0.38, 18, NAVY, SERIF, bold=True)
add_picture(slide, DATA["assets"]["precio_barrio_compacto"], 1.02, 2.34, w=5.05)
add_text(slide, "QUÉ OFRECE", 6.84, 1.66, 5.75, 0.3, 10, TERRA, bold=True)
add_text(slide, "Amenities y precio por m²", 6.84, 1.96, 5.75, 0.38, 18, NAVY, SERIF, bold=True)
add_picture(slide, DATA["assets"]["amenities_compacto"], 7.48, 2.34, w=4.35)
add_rect(slide, MX, 5.78, 5.85, 1.02, "ECE8E2")
add_text(slide, "Ubicación: señal principal", 0.88, 5.94, 2.8, 0.3, 15, NAVY, SERIF, bold=True)
add_text(slide, "Belgrano publica una mediana de USD 2.726/m², frente a USD 2.160/m² en el conjunto. El barrio muestra un efecto relevante sobre el precio.", 0.88, 6.29, 5.22, 0.43, 10.5, INK)
add_rect(slide, 6.84, 5.78, 5.75, 1.02, "ECE8E2")
add_text(slide, "Amenities: complemento", 7.08, 5.94, 3.25, 0.3, 14, NAVY, SERIF, bold=True)
add_text(slide, "La dotación alta se asocia con USD 392/m² adicionales, pero la relación es baja (rho = 0,19). No toda mejora recupera su costo.", 7.08, 6.27, 5.05, 0.46, 10, INK)

# 8. Insight: PCA y clusters
slide = prs.slides.add_slide(blank)
add_header(slide, "Insight 2 · Segmentación", 8, "CABA no es un único mercado: identificamos seis perfiles con lógicas distintas", title_size=25)
add_text(slide, "SIMPLIFICAR", MX, 1.66, 5.65, 0.3, 10, TERRA, bold=True)
add_text(slide, "Tres lentes resumen muchas variables", MX, 1.96, 5.65, 0.38, 18, NAVY, SERIF, bold=True)
add_picture(slide, DATA["assets"]["pca_negocio"], 1.08, 2.38, w=4.7)
add_text(slide, "SEGMENTAR", 6.58, 1.66, 5.95, 0.3, 10, TERRA, bold=True)
add_text(slide, "Seis perfiles de mercado", 6.58, 1.96, 5.95, 0.38, 18, NAVY, SERIF, bold=True)
add_picture(slide, DATA["assets"]["clusters_perfiles"], 6.93, 2.38, w=5.1)
add_rect(slide, MX, 5.78, 5.55, 1.02, "ECE8E2")
add_text(slide, "Qué aporta el PCA", 0.88, 5.94, 2.7, 0.3, 15, NAVY, SERIF, bold=True)
add_text(slide, "Convierte variables dispersas en indicadores comparables de amplitud, entorno urbano y confort. Facilita leer cada propiedad sin revisar decenas de columnas.", 0.88, 6.29, 4.95, 0.43, 10.5, INK)
add_rect(slide, 6.58, 5.78, 6.01, 1.02, "ECE8E2")
add_text(slide, "Decisión para el flipper", 6.82, 5.94, 3.0, 0.3, 15, NAVY, SERIF, bold=True)
add_text(slide, "Belgrano, Recoleta y Retiro forman el perfil de alto valor por m². Los comparables y las mejoras deben evaluarse dentro de ese mercado, no contra toda CABA.", 6.82, 6.29, 5.25, 0.43, 10.5, INK)

# 9. Segunda opinión
slide = prs.slides.add_slide(blank)
add_header(slide, "Paso 3 · Cómo contrastar", 9, "El modelo aporta una segunda opinión sobre el precio", title_size=28)
add_picture(slide, DATA["assets"]["importancia_modelo"], MX, 1.82, w=7.15)
add_text(slide, "¿Qué mira el modelo?", 8.32, 2.02, 3.9, 0.5, 23, TERRA, SERIF, bold=True)
add_bullets(slide, [
    "Superficie total y cubierta",
    "Expensas",
    "Ubicación",
    "Baños y antigüedad",
    "Estado y amenities",
], 8.32, 2.72, 3.75, 2.35, 13.5, INK, 7)
add_line(slide, 8.32, 5.24, 3.75)
add_text(slide, "No estima el precio de reventa.", 8.32, 5.48, 3.75, 0.42, 16, NAVY, SERIF, bold=True)
add_text(slide, "Sirve para detectar avisos cuyo precio publicado merece una revisión adicional.", 8.32, 5.98, 3.75, 0.72, 11.5, GRAY, italic=True)

# 10. Ejemplo
slide = prs.slides.add_slide(blank)
add_header(slide, "Los tres pasos en acción", 10, "Belgrano: una propiedad que reúne las señales para investigar", title_size=27)
prop = EXAMPLE["propiedad"]
rank = EXAMPLE["ranking"]
models = EXAMPLE["modelos"]
add_text(slide, "1 · PERFIL: ALTO VALOR/M²   2 · DESCUENTO + MEJORA   3 · SEGUNDA OPINIÓN", MX, 1.82, 7.1, 0.3, 9.5, TERRA, bold=True)
add_text(slide, f"{prop['barrio']} · {prop['calle']} {prop['altura']}", MX, 2.25, 5.7, 0.55, 22, NAVY, SERIF, bold=True)
add_text(slide, f"Departamento · {prop['ambientes']} ambientes · {prop['superficie_total_m2']:.0f} m² · {prop['estado'].lower()}", MX, 2.82, 5.75, 0.38, 12.5, GRAY)
facts = [
    ("Precio publicado", f"USD {rank['precio_publicado_usd']:,.0f}", f"USD {rank['precio_m2_usd']:,.0f}/m²"),
    ("Comparable", f"USD {rank['precio_m2_comparable_usd']:,.0f}/m²", f"{rank['comparables']} avisos equivalentes"),
    ("Subvaluación", f"{rank['subvaluacion_pct']:.1f}%", "descuento frente al comparable"),
    ("Índice de oportunidad", f"{rank['indice_oportunidad']:.0f}/100", "prioriza revisión; no calcula ROI"),
]
for i, (label, value, detail) in enumerate(facts):
    y = 3.35 + i * 0.76
    add_line(slide, MX, y - 0.08, 5.7)
    add_text(slide, label, MX, y, 2.35, 0.34, 11.5, GRAY)
    add_text(slide, value, 3.3, y - 0.03, 3.1, 0.42, 17, TERRA if i == 2 else NAVY, SERIF, bold=True, align=PP_ALIGN.RIGHT)
    add_text(slide, detail, MX, y + 0.37, 5.7, 0.25, 9.5, GRAY, italic=True)
add_rect(slide, 6.77, 1.92, 0.02, 4.75, LIGHT)
add_text(slide, "SEGUNDA OPINIÓN DEL MODELO", 7.2, 1.88, 5.2, 0.3, 10, TERRA, bold=True)
predictions = [
    ("Ridge", models["ridge"]["prediccion_usd"], models["ridge"]["diferencia_vs_publicado_pct"]),
    ("Random Forest", models["random_forest"]["prediccion_usd"], models["random_forest"]["diferencia_vs_publicado_pct"]),
]
for i, (name, value, diff) in enumerate(predictions):
    y = 2.48 + i * 0.78
    add_line(slide, 7.2, y - 0.08, 5.25)
    add_text(slide, name, 7.2, y, 2.1, 0.42, 13, NAVY)
    add_text(slide, f"USD {value:,.0f}", 9.12, y - 0.03, 2.2, 0.45, 18, NAVY, SERIF, bold=True, align=PP_ALIGN.RIGHT)
    add_text(slide, f"+{diff:.1f}%", 11.42, y, 1.0, 0.42, 13, TERRA, bold=True, align=PP_ALIGN.RIGHT)
add_text(slide, "Ambos modelos estiman un valor superior al publicado. Esa coincidencia justifica investigar el aviso, no comprarlo automáticamente.", 7.2, 4.18, 5.2, 1.05, 13.5, INK)
add_line(slide, 7.2, 5.45, 5.25)
add_text(slide, "ANTES DE INVERTIR", 7.2, 5.65, 3.0, 0.3, 10, NAVY, bold=True)
add_text(slide, "Validar estado real, precio negociado, obra, impuestos, comisiones y plazo de salida.", 7.2, 6.02, 5.2, 0.62, 12, INK)
add_text(slide, "Caso ilustrativo observado en train; no es validación fuera de muestra.", 7.2, 6.72, 5.2, 0.28, 9.5, GRAY, italic=True)

# 11. Impacto y decisiones
slide = prs.slides.add_slide(blank)
add_header(slide, "Impacto y recomendaciones", 11, "La herramienta concentra el esfuerzo donde hay más evidencia", title_size=27)
add_text(slide, "7.245", MX, 2.0, 2.0, 0.7, 34, NAVY, SERIF, bold=True)
add_text(slide, "propiedades comparables", MX, 2.75, 2.5, 0.5, 12, GRAY)
add_text(slide, "→", 2.7, 2.18, 0.55, 0.45, 21, TERRA, align=PP_ALIGN.CENTER)
add_text(slide, "3.129", 3.3, 2.0, 2.0, 0.7, 34, TERRA, SERIF, bold=True)
add_text(slide, "casos elegibles para investigar", 3.3, 2.75, 2.8, 0.5, 12, GRAY)
add_line(slide, MX, 3.52, 5.45)
add_text(slide, "Qué mejora para el negocio", MX, 3.78, 4.6, 0.4, 19, NAVY, SERIF, bold=True)
add_bullets(slide, [
    "Menos avisos revisados sin evidencia",
    "Comparables construidos con reglas consistentes",
    "Oportunidades ordenadas y explicables",
    "Validaciones pendientes visibles antes de invertir",
], MX, 4.35, 5.45, 1.85, 12.5, INK, 5)
metrics = DATA["metricas_modelos"]
add_text(slide, "MODELO COMO CONTROL ADICIONAL", 7.28, 1.9, 4.8, 0.3, 10, TERRA, bold=True)
headers = ["Modelo", "MAE", "MAPE", "R²"]
widths = [1.55, 1.25, 0.9, 0.65]
positions = [7.28, 8.83, 10.08, 10.98]
for x, width, header in zip(positions, widths, headers):
    add_rect(slide, x, 2.28, width, 0.43, NAVY)
    add_text(slide, header, x + 0.05, 2.37, width - 0.1, 0.22, 10.5, WHITE, bold=True, align=PP_ALIGN.CENTER)
for row_i, item in enumerate(metrics):
    y = 2.72 + row_i * 0.48
    values = [
        "Ridge" if row_i == 0 else "Random Forest",
        f"USD {item['MAE_USD']:,.0f}",
        f"{item['MAPE_%']:.1f}%",
        f"{item['R2']:.2f}",
    ]
    fill = "EAE0CF" if row_i == 1 else PAPER
    for x, width, value in zip(positions, widths, values):
        add_rect(slide, x, y, width, 0.46, fill, LIGHT)
        add_text(slide, value, x + 0.04, y + 0.1, width - 0.08, 0.23, 10.5, NAVY if row_i else INK, bold=row_i == 1, align=PP_ALIGN.CENTER)
add_text(slide, "Random Forest reduce el MAE 15,9% y mantiene mejor estabilidad en validación agrupada.", 7.28, 3.9, 4.4, 0.72, 13.5, NAVY, SERIF, bold=True)
add_line(slide, 7.28, 4.82, 4.38)
add_text(slide, "REGLAS PARA USARLO BIEN", 7.28, 5.02, 4.4, 0.3, 10, TERRA, bold=True)
add_bullets(slide, [
    "Concentrar el monitoreo en zonas con evidencia suficiente",
    "Priorizar descuentos defendibles, no scores aislados",
    "Usar el modelo como alerta de desalineamiento",
    "Cerrar siempre con due diligence financiera, técnica y legal",
], 7.28, 5.38, 4.9, 1.42, 11.5, INK, 3)

# 12. Futuro
slide = prs.slides.add_slide(blank)
add_header(slide, "Próximos pasos", 12, "El siguiente paso es pasar de priorización a rentabilidad medida", dark=True, title_size=27)
roadmap = [
    ("01", "Historia", "Guardar cambios de precio y días en mercado."),
    ("02", "Cierre", "Incorporar precios negociados y efectivos."),
    ("03", "Costos", "Registrar obra, impuestos, comisiones y financiación."),
    ("04", "Operación", "Automatizar scores, alertas y actualización de Power BI."),
    ("05", "Aprendizaje", "Monitorear resultados y reentrenar el modelo."),
]
for i, (num, title, body) in enumerate(roadmap):
    y = 1.93 + i * 0.88
    add_text(slide, num, MX, y, 0.55, 0.35, 17, TERRA, SERIF, bold=True)
    add_text(slide, title, 1.38, y, 1.35, 0.34, 14, WHITE, bold=True)
    add_text(slide, body, 2.82, y, 5.1, 0.42, 12.5, ICE)
    if i < 4:
        add_line(slide, MX, y + 0.56, 7.25, "2A3C52")
add_rect(slide, 8.75, 1.95, 0.02, 3.0, TERRA)
add_text(slide, "CUANDO EXISTAN LAS VARIABLES FINANCIERAS", 9.05, 2.0, 3.5, 0.52, 9.5, TERRA, bold=True)
add_text(slide, "ROI =", 9.05, 2.82, 3.2, 0.55, 28, WHITE, SERIF, bold=True)
add_text(slide, "(venta − compra − obra − gastos)\ncapital invertido", 9.05, 3.48, 3.45, 1.25, 16, ICE, SERIF, italic=True, align=PP_ALIGN.CENTER)
add_line(slide, MX, 6.12, W - 2 * MX, "2A3C52")
add_text(slide, "Hasta entonces, el producto maximiza la calidad de las oportunidades detectadas. No promete retornos.", MX, 6.35, 7.7, 0.55, 14, WHITE, SERIF, italic=True)
add_text(slide, "Precios publicados · Asociaciones no causales · Modelo como segunda opinión", 8.75, 6.42, 3.75, 0.35, 9, ICE, align=PP_ALIGN.RIGHT)

prs.save(OUTPUT)
print(f"OK -> {OUTPUT}")
